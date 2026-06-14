"""
File extractor for converting raw downloaded documents to plain text files.
"""
import os
import json
import shutil
import zipfile
import logging
from pathlib import Path
from typing import Dict, List, Optional
from datetime import datetime
import pandas as pd

# Try to import dependencies and handle gracefully if they're not available
try:
    from PyPDF2 import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False


logger = logging.getLogger(__name__)


class FileExtractor:
    """Convert raw downloaded files in data/raw/documents/{SYMBOL}/ into plain text files in data/trans/documents/{SYMBOL}/."""
    
    def __init__(self, symbol: str):
        self.symbol = symbol
        self.raw_dir = Path(f"data/raw/documents/{symbol}")
        self.trans_dir = Path(f"data/trans/documents/{symbol}")
        self.trans_dir.mkdir(parents=True, exist_ok=True)
        
    def extract_pdf(self, file_path: Path) -> tuple[str, int, int]:
        """
        Extract text from PDF files, handling encrypted PDFs gracefully.
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            Tuple of (extracted_filename, char_count, page_count)
        """
        filename = file_path.name
        extracted_filename = f"{Path(filename).stem}.txt"
        output_path = self.trans_dir / extracted_filename
        
        try:
            # Try to read the PDF with PyPDF2 if available
            if not PDF_AVAILABLE:
                raise ImportError("PyPDF2 is required for PDF extraction")
                
            pdf_reader = PdfReader(file_path)
            page_count = len(pdf_reader.pages)
            text_content = ""
            
            for i, page in enumerate(pdf_reader.pages):
                if i > 0:
                    text_content += "\n\n--- PAGE " + str(i+1) + " ---\n\n"
                page_text = page.extract_text()
                text_content += page_text
            
            char_count = len(text_content)
            
        except Exception as e:
            # Handle encrypted PDFs or other errors by writing empty file with warning
            logger.warning(f"Failed to extract PDF {filename}: {str(e)}")
            text_content = ""
            char_count = 0
            page_count = 0
            
            # Write an empty .txt file to maintain structure
            with open(output_path, 'w') as f:
                f.write("")
                
        else:
            # Write extracted content
            with open(output_path, 'w') as f:
                f.write(text_content)
                
        return (extracted_filename, char_count, page_count)
    
    def extract_html(self, file_path: Path) -> tuple[str, int]:
        """
        Extract text from HTML files using BeautifulSoup.
        
        Args:
            file_path: Path to the HTML file
            
        Returns:
            Tuple of (extracted_filename, char_count)
        """
        filename = file_path.name
        extracted_filename = f"{Path(filename).stem}.txt"
        output_path = self.trans_dir / extracted_filename
        
        try:
            if not BS4_AVAILABLE:
                raise ImportError("BeautifulSoup is required for HTML extraction")
                
            with open(file_path, 'r', encoding='utf-8') as f:
                soup = BeautifulSoup(f, 'html.parser')
                
            # Remove unwanted tags
            for tag in soup(['script', 'style', 'nav', 'header', 'footer', 'aside']):
                tag.decompose()
                
            text_content = soup.get_text()
            char_count = len(text_content)
            
        except Exception as e:
            logger.error(f"Failed to extract HTML {filename}: {str(e)}")
            text_content = ""
            char_count = 0
            
        # Write extracted content
        with open(output_path, 'w') as f:
            f.write(text_content)
            
        return (extracted_filename, char_count)
    
    def extract_zip(self, file_path: Path) -> tuple[str, int, int]:
        """
        Safely unzip files handling zip-slip attacks.
        
        Args:
            file_path: Path to the ZIP file
            
        Returns:
            Tuple of (extracted_filename, char_count, page_count)
        """
        filename = file_path.name
        extracted_filename = f"{Path(filename).stem}.txt"
        output_path = self.trans_dir / extracted_filename
        
        # Extract to temporary directory
        temp_extract_dir = self.trans_dir / f"{Path(filename).stem}_extracted"
        temp_extract_dir.mkdir(exist_ok=True)
        
        char_count = 0
        page_count = 0
        
        try:
            with zipfile.ZipFile(file_path, 'r') as zip_file:
                for member in zip_file.namelist():
                    # Prevent zip-slip vulnerability
                    member_path = Path(member)
                    if member_path.parts[0] == '..':
                        continue
                    
                    # Ensure path is inside our extraction directory
                    target_path = temp_extract_dir / member_path
                    if target_path.resolve().parent != temp_extract_dir.resolve():
                        continue
                        
                    # Extract file
                    if not target_path.is_dir():
                        with open(target_path, 'wb') as f:
                            f.write(zip_file.read(member))
                        
                        # If the extracted file is a supported type, process it recursively
                        if member_path.suffix.lower() in ['.pdf', '.html', '.txt', '.csv']:
                            sub_extractor = FileExtractor(self.symbol)
                            try:
                                if member_path.suffix.lower() == '.pdf':
                                    sub_filename, sub_chars, sub_pages = sub_extractor.extract_pdf(target_path)
                                elif member_path.suffix.lower() == '.html':
                                    sub_filename, sub_chars = sub_extractor.extract_html(target_path)
                                else:  # text-based files
                                    with open(target_path, 'r') as f:
                                        content = f.read()
                                    char_count += len(content)
                                    continue
                                    
                                char_count += sub_chars
                                if sub_pages and isinstance(sub_pages, int):
                                    page_count += sub_pages
                            except Exception:
                                pass  # Continue processing other files
                        else:
                            # For unsupported file types, we don't count their content
                            pass
        
        except Exception as e:
            logger.error(f"Failed to extract ZIP {filename}: {str(e)}")
            char_count = 0
            page_count = 0
            
        finally:
            # Clean up the temporary directory
            if temp_extract_dir.exists():
                shutil.rmtree(temp_extract_dir)
            
        # Write final content (or empty if failed)
        with open(output_path, 'w') as f:
            f.write("")
            
        return (extracted_filename, char_count, page_count)
    
    def extract_office(self, file_path: Path) -> tuple[str, int]:
        """
        Extract text from office files (.xlsx, .xls, .pptx, .docx).
        
        Args:
            file_path: Path to the office file
            
        Returns:
            Tuple of (extracted_filename, char_count)
        """
        filename = file_path.name
        extracted_filename = f"{Path(filename).stem}.txt"
        output_path = self.trans_dir / extracted_filename
        
        text_content = ""
        char_count = 0
        
        try:
            if file_path.suffix.lower() in ['.xlsx', '.xls']:
                # Extract from Excel files
                if not PDF_AVAILABLE:  # Using pandas for excel
                    raise ImportError("pandas required for Excel extraction")
                
                sheets = pd.read_excel(file_path, sheet_name=None)
                for sheet_name, df in sheets.items():
                    text_content += f"\n--- Sheet: {sheet_name} ---\n"
                    text_content += df.to_string(index=False)
                    text_content += "\n"
                    
            elif file_path.suffix.lower() == '.pptx' and PPTX_AVAILABLE:
                # Extract from PowerPoint files
                prs = Presentation(file_path)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content += shape.text + "\n"
                    
            elif file_path.suffix.lower() == '.docx' and DOCX_AVAILABLE:
                # Extract from Word files
                doc = Document(file_path)
                for para in doc.paragraphs:
                    text_content += para.text + "\n"
                    
            else:
                # Unsupported format, return empty content
                text_content = ""
                
            char_count = len(text_content)
            
        except Exception as e:
            logger.error(f"Failed to extract office file {filename}: {str(e)}")
            text_content = ""
            char_count = 0
            
        # Write extracted content
        with open(output_path, 'w') as f:
            f.write(text_content)
            
        return (extracted_filename, char_count)
    
    def extract_file(self, file_path: Path) -> tuple[str, str, int, int, bool, Optional[str], Optional[str]]:
        """
        Dispatch each file to the right extractor by extension.
        
        Args:
            file_path: Path to source file
            
        Returns:
            Tuple of (extracted_filename, extractor_type, char_count, page_count, success, error, source_url)
        """
        filename = file_path.name
        ext = file_path.suffix.lower()
        
        # Get the source URL from manifest if it exists
        source_url = None
        try:
            with open(self.raw_dir / "manifest.json", 'r') as f:
                manifest = json.load(f)
                # Try to find source_url in manifest - structure may vary
                for item in manifest.get('files', []):
                    if item.get('original') == filename:
                        source_url = item.get('source_url')
                        break
        except Exception:
            pass  # Continue without URL if not found
        
        try:
            if ext == '.pdf':
                extracted_filename, char_count, page_count = self.extract_pdf(file_path)
                return (extracted_filename, "pdf", char_count, page_count, True, None, source_url)
            
            elif ext in ['.html', '.htm']:
                extracted_filename, char_count = self.extract_html(file_path)
                return (extracted_filename, "html", char_count, 0, True, None, source_url)
                
            elif ext == '.zip':
                extracted_filename, char_count, page_count = self.extract_zip(file_path)
                return (extracted_filename, "zip", char_count, page_count, True, None, source_url)
                
            elif ext in ['.xlsx', '.xls']:
                extracted_filename, char_count = self.extract_office(file_path)
                return (extracted_filename, "xlsx", char_count, 0, True, None, source_url)
                
            elif ext in ['.pptx']:
                extracted_filename, char_count = self.extract_office(file_path)
                return (extracted_filename, "pptx", char_count, 0, True, None, source_url)
                
            elif ext in ['.docx']:
                extracted_filename, char_count = self.extract_office(file_path)
                return (extracted_filename, "docx", char_count, 0, True, None, source_url)
                
            elif ext == '.csv':
                # Copy CSV files as-is
                extracted_filename = filename
                output_path = self.trans_dir / extracted_filename
                shutil.copy2(file_path, output_path)
                return (extracted_filename, "csv", 0, 0, True, None, source_url)
                
            else:
                # Treat other file types as plain text files
                try:
                    with open(file_path, 'r') as f:
                        content = f.read()
                    char_count = len(content)
                    
                    extracted_filename = f"{Path(filename).stem}.txt"
                    output_path = self.trans_dir / extracted_filename
                    
                    with open(output_path, 'w') as f:
                        f.write(content)
                        
                    return (extracted_filename, "text", char_count, 0, True, None, source_url)
                    
                except Exception as e:
                    logger.error(f"Failed to copy/parse {filename}: {str(e)}")
                    extracted_filename = f"{Path(filename).stem}.txt"
                    output_path = self.trans_dir / extracted_filename
                    with open(output_path, 'w') as f:
                        f.write("")
                    return (extracted_filename, "unknown", 0, 0, False, str(e), source_url)
                    
        except Exception as e:
            logger.error(f"Failed to extract {filename}: {str(e)}")
            extracted_filename = f"{Path(filename).stem}.txt"
            return (extracted_filename, ext[1:] if ext else "unknown", 0, 0, False, str(e), source_url)
    
    def process_all_files(self) -> dict:
        """
        Process all files for this symbol and create index.json.
        
        Returns:
            Dictionary with symbol, extracted_at, and files as described in the schema
        """
        # Collect all files to process
        file_list = []
        if not self.raw_dir.exists():
            return {
                "symbol": self.symbol,
                "extracted_at": datetime.utcnow().isoformat() + "+00:00",
                "files": file_list
            }
        
        # Process each file in raw directory
        for file_path in self.raw_dir.iterdir():
            if file_path.is_file():
                extracted_filename, extractor, char_count, page_count, success, error, source_url = self.extract_file(file_path)
                file_list.append({
                    "original": file_path.name,
                    "extracted": extracted_filename,
                    "extractor": extractor,
                    "char_count": char_count,
                    "pages": page_count if extractor == "pdf" else None,
                    "success": success,
                    "error": error,
                    "source_url": source_url
                })
        
        # Create index.json file with the required schema
        index_data = {
            "symbol": self.symbol,
            "extracted_at": datetime.utcnow().isoformat() + "+00:00",
            "files": file_list
        }
        
        # Write the index file
        index_file = self.trans_dir / "index.json"
        with open(index_file, 'w') as f:
            json.dump(index_data, f, indent=2)
            
        return index_data