
"""Combined conversion functionality for data harvester pipeline."""

import os
import json
import shutil  
import zipfile
import logging
from pathlib import Path
from typing import Dict, List, Optional
from datetime import datetime
import pandas as pd

# Try to import dependencies and handle gracefully if they're not available
try:
    from PyPDF2 import PdfReader
    PDF_AVAILABLE = True
except ImportError:
    PDF_AVAILABLE = False

try:
    from bs4 import BeautifulSoup
    BS4_AVAILABLE = True
except ImportError:
    BS4_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from docx import Document
    DOCX_AVAILABLE = True
except ImportError:
    DOCX_AVAILABLE = False

logger = logging.getLogger(__name__)

class FileExtractor:
    """Convert raw downloaded files in data/raw/documents/{SYMBOL}/ into plain text files in data/trans/documents/{SYMBOL}/."""
    
    def __init__(self, symbol: str):
        self.symbol = symbol
        self.raw_dir = Path(f"data/raw/documents/{symbol}")
        self.trans_dir = Path(f"data/trans/documents/{symbol}")
        self.trans_dir.mkdir(parents=True, exist_ok=True)
        
    def extract_text_from_file(self, file_path: Path) -> Optional[str]:
        """Extract text from a file based on its extension."""
        if not file_path.exists():
            return None
            
        ext = file_path.suffix.lower()
        
        # Extract text based on file type  
        if ext == ".pdf" and PDF_AVAILABLE:
            try:
                reader = PdfReader(str(file_path))
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text
            except Exception as e:
                logger.error(f"Failed to extract PDF text: {e}")
                return None
                
        elif ext == ".docx" and DOCX_AVAILABLE:
            try:
                doc = Document(str(file_path))
                text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
                return text
            except Exception as e:
                logger.error(f"Failed to extract DOCX text: {e}")
                return None
                
        elif ext == ".pptx" and PPTX_AVAILABLE:
            try:
                presentation = Presentation(str(file_path))
                text = ""
                for slide in presentation.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            except Exception as e:
                logger.error(f"Failed to extract PPTX text: {e}")
                return None
                
        elif ext in (".txt", ".html", ".md") or BS4_AVAILABLE:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            except Exception as e:
                logger.error(f"Failed to extract text: {e}")
                # Try with different encoding  
                try:
                    with open(file_path, 'r', encoding='latin-1') as f:
                        return f.read()
                except Exception:
                    return None
                    
        else:
            # Generic file read (for unknown types or when dependencies unavailable) 
            try:
                with open(file_path, 'rb') as f:
                    return f.read().decode('utf-8', errors='ignore')
            except Exception:
                return None

    def extract_all(self):
        """Extract text from all files in raw directory and save to trans directory."""  
        logger.info(f"Extracting files for symbol {self.symbol}")
        
        if not self.raw_dir.exists():
            logger.info(f"No raw files found for {self.symbol}")
            return []

        extracted_files = []
        
        for file_path in self.raw_dir.iterdir():
            if file_path.is_file() and file_path.name.endswith(('.pdf', '.docx', '.pptx', '.txt', '.html', '.md')):
                try:
                    text_content = self.extract_text_from_file(file_path)
                    
                    if text_content:
                        # Generate output filename  
                        output_filename = f"{file_path.stem}.txt"
                        output_path = self.trans_dir / output_filename
                        
                        with open(output_path, 'w', encoding='utf-8') as f:
                            f.write(text_content)
                            
                        extracted_files.append({
                            "input_file": str(file_path),
                            "output_file": str(output_path),
                            "size": len(text_content), 
                            "success": True
                        })
                        
                except Exception as e:
                    logger.error(f"Failed to extract {file_path}: {e}")
                    extracted_files.append({
                        "input_file": str(file_path),
                        "error": str(e),
                        "success": False
                    })
                    
        return extracted_files

def flatten(src: Path, dest: Path) -> dict:
    """
    Read a raw numeric collector JSON from data/raw/ and write a normalized version to data/trans/numeric/.
    
    Args:
        src: Path to source raw JSON file
        dest: Path to destination flattened JSON file
        
    Returns:
        dict with symbol, source, rows, success, error keys
    """
    try:
        # Read input file  
        with open(src, 'r') as f:
            data = json.load(f)
        
        # Process and normalize the data 
        normalized_data = {}
        if isinstance(data, list):
            # Handle list of records (flatten each)
            for record in data:
                if "symbol" in record:
                    symbol = record["symbol"]
                    normalized_data[symbol] = record
                else:
                    # For simple records, use timestamp or index as identifier
                    normalized_data[f"record_{len(normalized_data)}"] = record
        else:
            # Single record
            normalized_data["0"] = data
            
        # Save normalized data  
        with open(dest, 'w') as f:
            json.dump(normalized_data, f, indent=2, ensure_ascii=False)
            
        return {
            "symbol": "UNKNOWN",  # This would be inferred from filename 
            "source": str(src),
            "rows": len(normalized_data),
            "success": True,
            "error": None
        }
    except Exception as e:
        logger.error(f"Error flattening {src}: {e}")
        return {
            "symbol": "UNKNOWN",
            "source": str(src),
            "rows": 0,
            "success": False,
            "error": str(e)
        }

